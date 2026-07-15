"""Generate a .pptx presentation for the fragment-translation-enhancement work.

Usage (inside the x2cangjie conda env, from repo root):
    pip install python-pptx
    python scripts/java/generate_presentation.py

Output: docs/presentation.pptx
"""
from __future__ import annotations

from pathlib import Path

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
except ImportError:
    raise SystemExit(
        "python-pptx is not installed. Run: pip install python-pptx"
    )


# ---------------------------------------------------------------------------
# Style constants
# ---------------------------------------------------------------------------
_BG = RGBColor(0x0D, 0x1B, 0x2A)        # dark navy
_ACCENT = RGBColor(0x00, 0xBE, 0xE6)    # cyan
_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
_LIGHT = RGBColor(0xD0, 0xD0, 0xD0)
_GREEN = RGBColor(0x4E, 0xC9, 0xB0)
_ORANGE = RGBColor(0xFF, 0xA5, 0x00)
_RED = RGBColor(0xFF, 0x55, 0x55)
_TABLE_HDR = RGBColor(0x1B, 0x2A, 0x3A)
_TABLE_ROW_ALT = RGBColor(0x12, 0x1F, 0x2E)

OUT = Path("docs/presentation.pptx")


# ---------------------------------------------------------------------------
# Helpers
# ---------------------------------------------------------------------------
def _set_bg(slide, color=_BG):
    """Fill slide background with a solid color."""
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_textbox(slide, left, top, width, height, anchor=MSO_ANCHOR.TOP):
    """Add a text box and return its text frame."""
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tf


def _set_para(para, text, size=18, color=_WHITE, bold=False, font_name="Calibri", alignment=PP_ALIGN.LEFT):
    """Configure a paragraph inside a text frame."""
    para.text = text
    para.font.size = Pt(size)
    para.font.color.rgb = color
    para.font.bold = bold
    para.font.name = font_name
    para.alignment = alignment


def _add_bullets(tf, items, size=16, color=_LIGHT, spacing=Pt(6)):
    """Add bullet paragraphs to a text frame."""
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_after = spacing
        if isinstance(item, tuple):
            text, sub_color, sub_bold = item
        else:
            text, sub_color, sub_bold = item, color, False
        _set_para(p, text, size=size, color=sub_color, bold=sub_bold)


def _add_title(slide, text, top=0.4, size=28, color=_ACCENT):
    """Add a slide title at the top."""
    tf = _add_textbox(slide, 0.6, top, 11.5, 0.8)
    _set_para(tf.paragraphs[0], text, size=size, color=color, bold=True)


def _add_table(slide, left, top, width, rows_data, col_widths=None, header=True):
    """Add a styled table. rows_data is a list of lists (first row = header if header=True)."""
    n_rows = len(rows_data)
    n_cols = len(rows_data[0])
    tbl_shape = slide.shapes.add_table(n_rows, n_cols, Inches(left), Inches(top), Inches(width), Inches(0.4 * n_rows))
    table = tbl_shape.table

    if col_widths:
        for i, w in enumerate(col_widths):
            table.columns[i].width = Inches(w)

    for r, row in enumerate(rows_data):
        for c, val in enumerate(row):
            cell = table.cell(r, c)
            cell.text = str(val)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(12)
            para.font.name = "Consolas" if r > 0 else "Calibri"
            if r == 0 and header:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _TABLE_HDR
                para.font.color.rgb = _ACCENT
                para.font.bold = True
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _TABLE_ROW_ALT if r % 2 == 0 else _BG
                para.font.color.rgb = _WHITE
    return tbl_shape


# ---------------------------------------------------------------------------
# Slide builders
# ---------------------------------------------------------------------------
def slide_cover(prs):
    """Slide 1: Cover."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank
    _set_bg(slide)

    tf = _add_textbox(slide, 1, 2.0, 10, 1.2, anchor=MSO_ANCHOR.MIDDLE)
    _set_para(tf.paragraphs[0], "Fragment 翻译增强", size=40, color=_WHITE, bold=True, alignment=PP_ALIGN.CENTER)
    p = tf.add_paragraph()
    _set_para(p, "原理、实现与消融分析", size=24, color=_ACCENT, alignment=PP_ALIGN.CENTER)

    tf2 = _add_textbox(slide, 1, 4.2, 10, 1.0, anchor=MSO_ANCHOR.MIDDLE)
    _set_para(tf2.paragraphs[0], "基于伪代码中间层、语法注入与结构检索的 LLM 代码翻译改进", size=16, color=_LIGHT, alignment=PP_ALIGN.CENTER)
    p2 = tf2.add_paragraph()
    _set_para(p2, "项目：x2cangjie  |  Java → 仓颉（Cangjie）自动翻译  |  2026-07-08", size=14, color=_LIGHT, alignment=PP_ALIGN.CENTER)


def slide_background(prs):
    """Slide 2: Background."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "背景：x2cangjie 在做什么")

    tf = _add_textbox(slide, 0.6, 1.2, 11.5, 1.0)
    _set_para(tf.paragraphs[0], "目标：把 Java 库自动翻译成仓颉（Cangjie）语言", size=18, color=_WHITE, bold=True)
    p = tf.add_paragraph()
    _set_para(p, "Pipeline: preprocess → create_schema → get_dependencies → translate_types → create_skeleton → build_mock_corpus → translate_fragment → analyze_errors", size=12, color=_LIGHT)
    p2 = tf.add_paragraph()
    _set_para(p2, "核心是第 7 步 translate_fragment：LLM 逐个 fragment 翻译，每填一个跑 cjpm build 编译验证，失败带错误反馈重试", size=14, color=_LIGHT)

    tf2 = _add_textbox(slide, 0.6, 3.0, 11.5, 0.5)
    _set_para(tf2.paragraphs[0], "观察到的两类系统性错误：", size=16, color=_ACCENT, bold=True)

    _add_table(slide, 0.6, 3.5, 11.0, [
        ["错误类型", "表现", "根因"],
        ["A: 错误继承 Java 源语法/API", "LLM 照搬 stream API / checked exception / for-each lambda", "LLM 模仿源码结构而非理解意图后用目标语言惯用法重写"],
        ["B: 使用错误的 Cangjie 语法/API", "泛型用 extends 而非 where T<:Bound、Any 当 HashMap key、boolean 而非 Bool", "Cangjie 是新语言，LLM 训练数据中几乎没有 Cangjie 代码"],
    ], col_widths=[3.0, 4.5, 3.5])

    tf3 = _add_textbox(slide, 0.6, 5.5, 11.5, 0.5)
    _set_para(tf3.paragraphs[0], "针对 A/B 两类错误，做了三件事（Part 1/2/3），互相独立、可单独或组合开关", size=16, color=_GREEN, bold=True)


def slide_part1(prs):
    """Slide 3: Part 1 — Pseudocode."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "Part 1：伪代码中间层（解决错误 A）")

    _add_table(slide, 0.6, 1.1, 11.0, [
        ["论文", "出处", "核心思想"],
        ["A1: Pseudocode-based Code Translation", "arXiv:2510.00920, 2025", "源→伪代码→目标，两阶段语义翻译"],
        ["A3: NL in the Middle", "CASCON 2025", "自然语言中间表示效果最好（+13.8%）"],
        ["A6: Assessing Intermediate Languages", "arXiv:2407.05411", "警示：收益可能部分来自 CoT 多步推理效应"],
    ], col_widths=[4.0, 3.0, 4.0])

    tf = _add_textbox(slide, 0.6, 3.2, 11.5, 0.5)
    _set_para(tf.paragraphs[0], "论文原理（A1）：", size=16, color=_ACCENT, bold=True)

    _add_bullets(_add_textbox(slide, 0.6, 3.6, 11.5, 1.5), [
        "直接翻译时 LLM 试图模仿源码结构 → 生成语义不一致的代码",
        "伪代码中间层：LLM 先把源码抽象成语言无关伪代码，再从伪代码生成目标代码",
        ("测试 5 种策略组合，发现 "伪代码 + 源代码"组合效果最好 — 伪代码居中解决歧义，源码作为 fallback", _GREEN, True),
    ], size=14)

    tf2 = _add_textbox(slide, 0.6, 5.3, 11.5, 0.5)
    _set_para(tf2.paragraphs[0], "项目实现：", size=16, color=_ACCENT, bold=True)

    _add_bullets(_add_textbox(slide, 0.6, 5.7, 11.5, 1.5), [
        "Java fragment → LLM → 伪代码+注释 (Phase-1) → 伪代码+Java源码+metadata → LLM → Cangjie (Phase-2)",
        "Phase-1 prompt 规约：仅通用关键字、API 调用改写为动词短语、每块前 // 注释说明意图",
        "失败退化为直接翻译；_skip_prompt_build 优化跳过昂贵上下文加载",
    ], size=14)


def slide_part2(prs):
    """Slide 4: Part 2 — Grammar EBNF."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "Part 2：Cangjie 语法 EBNF 注入（解决错误 B）")

    _add_table(slide, 0.6, 1.1, 11.0, [
        ["论文", "出处", "核心思想"],
        ["B2: Grammar Prompting", "Wang et al., ACL 2023", "BNF 语法注入 prompt，仅注入不需约束解码就显著提升"],
        ["B1: DocCGen", "EMNLP 2024", "从文档提取 grammar/schema 做约束解码，OOD 场景效果显著"],
    ], col_widths=[3.5, 3.0, 4.5])

    tf = _add_textbox(slide, 0.6, 2.6, 11.5, 0.5)
    _set_para(tf.paragraphs[0], "Cangjie 对 LLM 就是 DSL — 训练数据极少", size=16, color=_ORANGE, bold=True)

    tf2 = _add_textbox(slide, 0.6, 3.1, 11.5, 0.5)
    _set_para(tf2.paragraphs[0], "项目实现 — 注入两部分文本：", size=16, color=_ACCENT, bold=True)

    _add_table(slide, 0.6, 3.5, 11.0, [
        ["约束", "对应典型错误"],
        ["G1: 泛型用 where T <: Bound（不是 extends）", "Java ? extends T"],
        ["G3: Any 不满足 Hashable，用 AnyHashable", "HashMap<Object,V> 编译报错"],
        ["G5: 布尔类型是 Bool（不是 boolean）", "Java boolean"],
        ["G6: 字符串插值 ${expr}", "Java String.format"],
    ], col_widths=[6.0, 5.0])

    _add_bullets(_add_textbox(slide, 0.6, 5.8, 11.5, 1.0), [
        "第二部分：运行时 API 映射表（Object→AnyHashable, Runnable→()->Unit 等）",
        ("设计：规则在 YAML 配置可编辑；单例缓存；不做约束解码靠编译错误反馈做 rejection sampling", _LIGHT, False),
    ], size=14)


def slide_part3(prs):
    """Slide 5: Part 3 — Syntax-graph RAG."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "Part 3：语法图 RAG（CFG/DFG 结构相似检索）")

    _add_table(slide, 0.6, 1.1, 11.0, [
        ["论文", "出处", "核心思想"],
        ["B3: CodeGRAG", "Huang et al., arXiv:2405.02355, 2024", "提取 CFG+DFG 融合图，GNN+跨语言检索结构相似代码"],
        ["B4: Syntax-Aware RAG", "EMNLP 2023 Findings", "在 RAG 中引入语法感知，不只用语义相似还用语法结构相似度"],
    ], col_widths=[3.5, 3.5, 4.0])

    tf = _add_textbox(slide, 0.6, 2.6, 11.5, 0.5)
    _set_para(tf.paragraphs[0], "项目实现（实用化简化）— 纯正则结构指纹 + Jaccard 相似度，无 NN/CUDA/额外依赖", size=14, color=_ACCENT, bold=True)

    _add_table(slide, 0.6, 3.1, 11.0, [
        ["维度", "内容"],
        ["shape_bag", "12 个操作类别计数（cf_if/cf_loop/op_call/op_index 等），桶化为 0-3"],
        ["call_names", "方法调用点标识符集合"],
        ["container_types", "命中集合类型名（list/array/map/set 等）"],
    ], col_widths=[2.5, 8.5])

    _add_bullets(_add_textbox(slide, 0.6, 5.0, 11.5, 1.5), [
        ("检索：加权 Jaccard = 0.6×shape + 0.25×call + 0.15×container，返回 top-3", _WHITE, False),
        ("索引：扫描 CangjieCorpus，12874 个代码块，pickle 序列化", _WHITE, False),
        ("与现有 RAG 互补：原 RAG 回答"该用什么 API"，Part 3 回答"该写什么样的代码骨架"", _GREEN, True),
    ], size=14)


def slide_combination(prs):
    """Slide 6: How the three parts combine."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "三部分如何组合工作")

    tf = _add_textbox(slide, 0.6, 1.2, 11.5, 0.5)
    _set_para(tf.paragraphs[0], "Prompt 注入顺序：", size=16, color=_ACCENT, bold=True)

    tf2 = _add_textbox(slide, 0.6, 1.6, 11.5, 0.8)
    _set_para(tf2.paragraphs[0],
        "persona → instruction → grammar(Part2) → Java source → pseudocode(Part1) → partial → generics → KB → RAG docs → syntax_graph(Part3) → ICL → feedback → ### Response:",
        size=12, color=_LIGHT)

    tf3 = _add_textbox(slide, 0.6, 2.6, 11.5, 0.5)
    _set_para(tf3.paragraphs[0], "CLI 开关（三个 flag 默认 false，向后兼容）：", size=16, color=_ACCENT, bold=True)

    _add_table(slide, 0.6, 3.1, 11.0, [
        ["场景", "use_pseudocode", "use_grammar_prompt", "use_syntax_rag"],
        ["仅修 Java→Cangjie API 模式继承错", "true", "false", "false"],
        ["不熟悉 Cangjie 语法（多为编译报语法错）", "false", "true", "false"],
        ["需要 few-shot 结构模板", "false", "false", "true"],
        ["全开（增益最高）", "true", "true", "true"],
    ], col_widths=[5.0, 2.0, 2.0, 2.0])

    _add_bullets(_add_textbox(slide, 0.6, 5.5, 11.5, 0.8), [
        ("设计逻辑：grammar 在最前（先读规则再读代码）；伪代码在源码后（理解意图后再翻译）", _LIGHT, False),
        ("结构示例在 RAG 文档后、ICL 前（作为"怎么写"的模板参考）", _LIGHT, False),
    ], size=14)


def slide_code_structure(prs):
    """Slide 7: Code structure."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "代码结构")

    tf = _add_textbox(slide, 0.6, 1.1, 5.5, 3.0)
    _set_para(tf.paragraphs[0], "新增文件：", size=16, color=_GREEN, bold=True)
    _add_bullets(tf, [
        "grammar_prompt.py (Part 2)",
        "syntax_graph.py (Part 3)",
        "ablation_compare.py (消融分析)",
        "build_syntax_graph_index.sh",
        "run_ablation.sh",
        "test_grammar_prompt.py",
        "test_syntax_graph.py",
        "test_ablation_compare.py",
    ], size=13)

    tf2 = _add_textbox(slide, 6.5, 1.1, 5.5, 3.0)
    _set_para(tf2.paragraphs[0], "修改文件：", size=16, color=_ORANGE, bold=True)
    _add_bullets(tf2, [
        "compositional_translation_validation.py",
        "prompt_generator.py (Part 1/2/3 注入点)",
        "prompt_templates.yaml (新模板)",
        "translate_fragment.sh (3 个新参数)",
    ], size=13)

    _add_table(slide, 0.6, 4.3, 11.0, [
        ["想了解", "读这个文件"],
        ["翻译主循环", "compositional_translation_validation.py → translate()"],
        ["prompt 组装", "prompt_generator.py → build_base_prompt()"],
        ["Part 1 伪代码生成", "compositional_translation_validation.py → _generate_pseudocode()"],
        ["Part 2 语法规则", "configs/prompt_templates.yaml → cangjie_grammar_context"],
        ["Part 3 结构指纹", "syntax_graph.py → infer_structural_signature()"],
        ["消融报告生成", "ablation_compare.py → main()"],
    ], col_widths=[4.0, 7.0])


def slide_ablation_design(prs):
    """Slide 8: Ablation design."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "消融实验设计")

    tf = _add_textbox(slide, 0.6, 1.1, 11.5, 0.8)
    _set_para(tf.paragraphs[0], "动机（A6 论文警示）：", size=16, color=_ACCENT, bold=True)
    p = tf.add_paragraph()
    _set_para(p, "伪代码中间层的收益可能部分来自 CoT 多步推理效应，需要 ablation 分离每部分的增量来源", size=14, color=_LIGHT)

    tf2 = _add_textbox(slide, 0.6, 2.2, 11.5, 0.5)
    _set_para(tf2.paragraphs[0], "8 种 run-tag（2³ = 8）：", size=16, color=_ACCENT, bold=True)

    _add_table(slide, 0.6, 2.7, 11.0, [
        ["run-tag", "use_pseudocode", "use_grammar_prompt", "use_syntax_rag"],
        ["baseline", "false", "false", "false"],
        ["pseudo", "true", "false", "false"],
        ["grammar", "false", "true", "false"],
        ["syntax", "false", "false", "true"],
        ["pseudo+grammar", "true", "true", "false"],
        ["pseudo+syntax", "true", "false", "true"],
        ["grammar+syntax", "false", "true", "true"],
        ["all", "true", "true", "true"],
    ], col_widths=[3.0, 2.7, 2.7, 2.6])

    _add_bullets(_add_textbox(slide, 0.6, 6.2, 11.5, 0.8), [
        "配置：commons-csv / 381 fragments / gpt-4o / 温度 0.0",
        "每组前重建 skeleton；每组后 snapshot schema；Fisher exact 双侧 p 值（纯 Python）",
    ], size=13)


def slide_ablation_results(prs):
    """Slide 9: Ablation results."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "消融结果（commons-csv / gpt-4o / 381 fragments）")

    _add_table(slide, 0.6, 1.1, 11.0, [
        ["Run tag", "完成数", "完成率", "Δ vs baseline"],
        ["baseline", "241", "63.3%", "—"],
        ["pseudo (Part 1)", "249", "65.4%", "+2.1pp"],
        ["grammar (Part 2)", "255", "66.9%", "+3.7pp"],
        ["syntax (Part 3)", "258", "67.7%", "+4.5pp"],
        ["pseudo+grammar (1+2)", "259", "68.0%", "+4.7pp"],
        ["pseudo+syntax (1+3)", "260", "68.2%", "+5.0pp"],
        ["grammar+syntax (2+3)", "260", "68.2%", "+5.0pp"],
        ["all (1+2+3)", "260", "68.2%", "+5.0pp"],
    ], col_widths=[3.5, 2.0, 2.0, 3.5])

    _add_bullets(_add_textbox(slide, 0.6, 5.8, 11.5, 1.5), [
        ("单部分独立效应排序：Part 3 (+4.5pp) > Part 2 (+3.7pp) > Part 1 (+2.1pp)", _GREEN, True),
        ("组合饱和：两两组合接近 +5.0pp，三部分全开无额外增益 — 覆盖的错误类型有重叠", _ORANGE, True),
        ("显著性：Fisher exact p 值均 >0.05（单项目 381 样本量不足），但趋势清晰一致", _LIGHT, False),
    ], size=14)


def slide_analysis(prs):
    """Slide 10: Results analysis & future work."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "结果分析与未来工作")

    tf = _add_textbox(slide, 0.6, 1.1, 11.5, 0.5)
    _set_para(tf.paragraphs[0], "关键发现：", size=16, color=_ACCENT, bold=True)

    _add_bullets(_add_textbox(slide, 0.6, 1.5, 11.5, 3.0), [
        ("Part 3 语法图 RAG 单独贡献最大（+4.5pp）— 结构相似的 Cangjie 代码片段是最有效的 few-shot 示例", _WHITE, False),
        ("Part 2 语法注入次之（+3.7pp）— EBNF 规则直接减少语法类编译错误", _WHITE, False),
        ("Part 1 伪代码贡献最小（+2.1pp）— CoT 效应被 Part 2/3 部分吸收", _WHITE, False),
        ("组合饱和：Part 2+3 或 1+2+3 都达到 +5.0pp，覆盖错误类型有重叠", _WHITE, False),
        ("代价：Part 1 增加每 fragment 耗时 ~30%（多一次 LLM 调用），Part 2/3 几乎无额外开销", _WHITE, False),
    ], size=14)

    tf2 = _add_textbox(slide, 0.6, 4.8, 11.5, 0.5)
    _set_para(tf2.paragraphs[0], "后续可改进：", size=16, color=_ACCENT, bold=True)

    _add_bullets(_add_textbox(slide, 0.6, 5.2, 11.5, 1.5), [
        "扩展到多项目多模型（jansi / commons-cli + deepseek-chat / glm-5.1），增大样本量",
        "跑 mock 测试拿 test_pass 指标",
        "Part 1 做 CoT-only ablation 分离 CoT 效应",
        "Part 3 升级为 tree-sitter 真实 CFG/DFG + 跨语言预训练模型",
    ], size=14)


def slide_references(prs):
    """Slide 11: References."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _set_bg(slide)
    _add_title(slide, "参考论文速查")

    _add_table(slide, 0.6, 1.2, 11.0, [
        ["ID", "论文", "对应 Part"],
        ["A1", "Pseudocode-based Code Translation (arXiv 2510.00920)", "Part 1"],
        ["A3", "NL in the Middle (CASCON 2025)", "Part 1 设计依据"],
        ["A6", "Assessing Intermediate Languages (arXiv 2407.05411)", "Part 1 ablation 依据"],
        ["B2", "Grammar Prompting (ACL 2023)", "Part 2"],
        ["B1", "DocCGen (EMNLP 2024)", "Part 2 补充"],
        ["B3", "CodeGRAG (arXiv 2405.02355)", "Part 3"],
        ["B4", "Syntax-Aware RAG (EMNLP 2023 Findings)", "Part 3 补充"],
    ], col_widths=[1.0, 7.5, 2.5])

    _add_bullets(_add_textbox(slide, 0.6, 5.0, 11.5, 1.5), [
        "完整论文摘要：docs/related_work_code_translation.md",
        "完整实现细节：docs/fragment_translation_enhancements.md",
        "完整工作汇报：docs/work_report.md",
    ], size=14, color=_LIGHT)


# ---------------------------------------------------------------------------
# Main
# ---------------------------------------------------------------------------
def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide_cover(prs)
    slide_background(prs)
    slide_part1(prs)
    slide_part2(prs)
    slide_part3(prs)
    slide_combination(prs)
    slide_code_structure(prs)
    slide_ablation_design(prs)
    slide_ablation_results(prs)
    slide_analysis(prs)
    slide_references(prs)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))
    print(f"Generated: {OUT}  ({len(prs.slides._sldIdLst)} slides)")
    print(f"Path: {OUT.resolve()}")


if __name__ == "__main__":
    main()
